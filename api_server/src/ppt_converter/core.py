"""PPT(.pptx) → JSON 변환 핵심 로직.

설계 원칙:
- 슬라이드 1장 = 1 섹션 (level 1; 제목에 "1.2" 같은 번호 패턴이 있으면 level 2/3).
- 슬라이드 내 도형 등장 순서 = blocks 배열 순서.
- 표/그림은 본문 흐름에 ``ref`` 블록으로 삽입되고 데이터는 최상위 ``tables`` /
  ``figures`` 에 저장된다.
- 모든 그림은 ``attachments`` 배열에 ``kind=figure`` 로도 등록된다 (Word 변환기와 동일).
- 슬라이드 노트는 ``[Speaker Notes]`` 마커 단락 뒤에 본문 흐름으로 추가된다.
- 시각적 좌표(왼쪽/오른쪽 텍스트박스 등)는 보존되지 않는다 — reading order =
  python-pptx shape iteration order.
"""
from __future__ import annotations

import hashlib
import logging
import re
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.presentation import Presentation as PresentationType

from .models import (
    Attachment,
    Block,
    ConversionResult,
    Figure,
    Section,
    Source,
    Table,
)
from .parser import (
    extract_chart_title,
    extract_picture_alt_text,
    extract_picture_blob,
    extract_slide_title,
    extract_speaker_notes,
    extract_table_data,
    extract_text_lines,
    infer_body_heading,
    infer_section_id_from_title,
    is_chart_shape,
    is_picture_shape,
    is_table_shape,
    is_text_shape,
    iter_body_shapes,
    iter_slides_with_index,
)

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Migration 0007 — agent-discovery 자동 기본값 헬퍼
# ---------------------------------------------------------------------------

def _default_agent_hints(
    *, data_type_name: str, tags: list[str], section_count: int,
    table_count: int, figure_count: int,
) -> str:
    topics = ", ".join(t for t in (tags or []) if t) or "N/A"
    return (
        f"이 record 는 {data_type_name} 입니다. 주요 토픽: {topics}. "
        f"본문은 총 {section_count} 섹션, {table_count} 표, "
        f"{figure_count} 그림으로 구성됩니다."
    )


def _default_query_examples(*, title: str, tags: list[str]) -> list[str]:
    out: list[str] = []
    t = (title or "").strip()
    if t:
        out.append(f"{t} 어떻게 사용해?")
    first_tag = next((x for x in (tags or []) if isinstance(x, str) and x.strip()), None)
    if first_tag:
        out.append(f"{first_tag}에 대해 알려줘")
    if t:
        out.append(f"{t} 관련 자료 보여줘")
    return out[:3]


def _apply_agent_discovery_defaults(
    meta: dict[str, Any],
    *,
    overrides: dict[str, Any] | None = None,
    data_type_name: str,
    title: str,
    tags: list[str],
    section_count: int,
    table_count: int,
    figure_count: int,
) -> None:
    overrides = overrides or {}
    if "agent_hints" in overrides and overrides["agent_hints"] is not None:
        meta["agent_hints"] = overrides["agent_hints"]
    elif meta.get("agent_hints") is None:
        meta["agent_hints"] = _default_agent_hints(
            data_type_name=data_type_name, tags=tags,
            section_count=section_count, table_count=table_count,
            figure_count=figure_count,
        )
    if "related_record_ids" in overrides and overrides["related_record_ids"] is not None:
        meta["related_record_ids"] = list(overrides["related_record_ids"])
    elif "related_record_ids" not in meta:
        meta["related_record_ids"] = []
    if "query_examples" in overrides and overrides["query_examples"] is not None:
        meta["query_examples"] = list(overrides["query_examples"])
    elif not meta.get("query_examples"):
        meta["query_examples"] = _default_query_examples(title=title, tags=tags)
    if "access_pattern" in overrides and overrides["access_pattern"]:
        meta["access_pattern"] = overrides["access_pattern"]
    elif not meta.get("access_pattern"):
        meta["access_pattern"] = "occasional"


# ---------------------------------------------------------------------------
# Options
# ---------------------------------------------------------------------------
@dataclass
class PptxConverterOptions:
    """PPT 변환기 옵션."""

    division: str
    team: str
    year: int
    seq: int = 1
    output_dir: Path = field(default_factory=lambda: Path("output"))
    extract_images: bool = True
    tags: list[str] = field(default_factory=list)
    agents: list[str] = field(default_factory=list)
    # 실 데이터 적응 휴리스틱 (default ON — 작성자가 양식을 안 지킨 PPT 도 RAG-친화).
    extract_summary: bool = True              # core.subject / 본문에서 summary 폴백 추출
    group_consecutive_duplicates: bool = True  # 연속 동일 제목을 부모-자식 트리로 그룹화
    extract_body_headings: bool = True         # 본문 안의 "1.1 …" 패턴을 sub-section 으로 승격
    infer_caption_from_neighbor: bool = True   # 그림 alt-text/인접 텍스트박스로 캡션 자동 추정


# 캡션 접두어 패턴 — 이미 "Figure N:" / "Fig.N" / "그림 N:" 형태이면 그대로 사용.
_CAPTION_PREFIX_PATTERN = re.compile(
    r"^\s*(figure|fig\.?|그림|figura|图|圖|画像|abbildung)\s*\d+\s*[:\.]",
    re.IGNORECASE,
)

# 캡션 후보 최대 길이 (문자 수). 너무 긴 텍스트박스는 본문 단락이지 캡션이 아님.
_CAPTION_MAX_LEN = 120


# ---------------------------------------------------------------------------
# ID helpers (Word 변환기와 동일한 포맷)
# ---------------------------------------------------------------------------
def _make_doc_id(opts: PptxConverterOptions) -> str:
    return f"DOC-{opts.division}-{opts.team}-{opts.year}-{opts.seq:06d}"


def _make_fig_id(doc_id: str, n: int) -> str:
    return f"{doc_id}-F{n:03d}"


def _make_tbl_id(doc_id: str, n: int) -> str:
    return f"{doc_id}-T{n:03d}"


def _make_att_id(doc_id: str, n: int) -> str:
    return f"{doc_id}-A{n:03d}"


def _sha256_of_file(path: Path | str, chunk: int = 65536) -> str:
    h = hashlib.sha256()
    try:
        with open(path, "rb") as f:
            for buf in iter(lambda: f.read(chunk), b""):
                h.update(buf)
        return h.hexdigest()
    except OSError:
        return ""


# ---------------------------------------------------------------------------
# Converter
# ---------------------------------------------------------------------------
class PptxConverter:
    """PowerPoint 프레젠테이션을 DOC type JSON 으로 변환."""

    def __init__(self, options: PptxConverterOptions) -> None:
        self.opts = options
        self.doc_id = _make_doc_id(options)
        self.warnings: list[str] = []

        self.section_root: list[Section] = []
        self.figures: list[Figure] = []
        self.tables: list[Table] = []
        self.sources: list[Source] = []
        self.attachments: list[Attachment] = []

        self.fig_counter = 0
        self.tbl_counter = 0
        self.att_counter = 0

        # 자동 섹션 번호 (제목에 번호가 없는 슬라이드용)
        self.l1_counter = 0
        self.l2_counter = 0
        self.l3_counter = 0

    # ---- public API --------------------------------------------------

    def convert(self, pptx_path: str | Path) -> ConversionResult:
        pptx_path = Path(pptx_path)
        prs = Presentation(str(pptx_path))

        for sinfo in iter_slides_with_index(prs):
            self._process_slide(sinfo.index, sinfo.slide)

        # 실 데이터 적응 후처리: 본문 안의 "1.1 …" 패턴을 sub-section 으로 승격.
        # group_consecutive_duplicates 보다 먼저 실행해야 슬라이드별 children 이 정상 부착된다.
        if self.opts.extract_body_headings:
            self._extract_body_subsections()

        # 실 데이터 적응 후처리: 연속 동일 제목을 부모-자식 트리로 그룹화.
        if self.opts.group_consecutive_duplicates:
            self._group_consecutive_duplicate_titles()

        meta = self._build_meta(prs, pptx_path)
        return ConversionResult(
            schema_version="1.0",
            meta=meta,
            sections=self.section_root,
            figures=self.figures,
            tables=self.tables,
            sources=self.sources,
            attachments=self.attachments,
            warnings=self.warnings,
        )

    # ---- slide processing -------------------------------------------

    def _process_slide(self, slide_idx: int, slide: Any) -> None:
        title_text = extract_slide_title(slide)
        section = self._open_section_for_slide(slide_idx, title_text)

        # 본문 도형 순회.
        for shape in iter_body_shapes(slide):
            if is_table_shape(shape):
                self._handle_table_shape(section, shape)
            elif is_picture_shape(shape):
                self._handle_picture_shape(section, shape)
            elif is_chart_shape(shape):
                self._handle_chart_shape(section, shape)
            elif is_text_shape(shape):
                self._handle_text_shape(section, shape)
            # 그 외 (자동도형 안의 텍스트, 커넥터 등) — 텍스트 프레임이 있으면 위에서 잡힘.

        # 슬라이드 노트.
        notes = extract_speaker_notes(slide)
        if notes:
            section.blocks.append(Block(type="paragraph", text="[Speaker Notes]"))
            for line in notes.splitlines():
                line = line.strip()
                if not line:
                    continue
                section.blocks.append(Block(type="paragraph", text=line))

        # 빈 슬라이드 경고
        if not section.blocks and not section.title:
            self.warnings.append(
                f"슬라이드 {slide_idx}: 제목/본문/노트가 모두 비어 있음"
            )

    # ---- section management -----------------------------------------

    def _open_section_for_slide(self, slide_idx: int, title_text: str) -> Section:
        """슬라이드 제목으로부터 섹션 1개를 만들어 트리에 등록.

        제목이 "1.2 작동원리" 형태이면 level=2 로 부모 섹션 찾기 시도, 부모가
        없으면 가상 부모를 만들지 않고 level=1 로 폴백한다 (단순화).
        제목이 비어 있으면 자동 번호 + "(슬라이드 N)" 제목 사용.
        """
        if not title_text:
            self.warnings.append(
                f"슬라이드 {slide_idx}: 제목 placeholder 누락 — 자동 제목 사용"
            )
            title = f"슬라이드 {slide_idx}"
            section_id = self._next_auto_id(level=1)
            section = Section(id=section_id, level=1, title=title)
            self.section_root.append(section)
            return section

        parsed_id, clean_title = infer_section_id_from_title(title_text)

        if parsed_id is not None:
            level = parsed_id.count(".") + 1
            section = Section(id=parsed_id, level=level, title=clean_title)
            self._sync_counters(parsed_id)
            self._attach_section(section, level)
            return section

        # 번호 없는 제목 → level 1 자동 번호.
        section_id = self._next_auto_id(level=1)
        section = Section(id=section_id, level=1, title=clean_title)
        self.section_root.append(section)
        return section

    def _attach_section(self, section: Section, level: int) -> None:
        """level 2/3 섹션을 가장 가까운 상위 레벨 섹션의 children 으로 붙인다.

        매칭되는 부모가 없으면 level 1 로 폴백하여 root 에 추가하고 경고.
        """
        if level == 1:
            self.section_root.append(section)
            return

        parent = self._find_last_section_at_level(level - 1)
        if parent is None:
            # 부모 없음 → root 로 폴백 + 경고.
            self.warnings.append(
                f"섹션 {section.id} (level {level}): 상위 레벨 섹션이 없어 root 로 등록"
            )
            section.level = 1
            self.section_root.append(section)
            return
        parent.children.append(section)

    def _find_last_section_at_level(self, level: int) -> Section | None:
        """현재까지 등록된 섹션 중 지정 level 의 마지막 섹션을 반환."""
        candidates: list[Section] = []

        def walk(nodes: list[Section]) -> None:
            for n in nodes:
                if n.level == level:
                    candidates.append(n)
                walk(n.children)

        walk(self.section_root)
        return candidates[-1] if candidates else None

    def _next_auto_id(self, level: int) -> str:
        if level == 1:
            self.l1_counter += 1
            self.l2_counter = 0
            self.l3_counter = 0
            return f"{self.l1_counter}"
        if level == 2:
            self.l2_counter += 1
            self.l3_counter = 0
            return f"{self.l1_counter}.{self.l2_counter}"
        self.l3_counter += 1
        return f"{self.l1_counter}.{self.l2_counter}.{self.l3_counter}"

    def _sync_counters(self, section_id: str) -> None:
        """원본 제목의 번호와 자동 카운터를 맞춰둔다 (이후 자동번호용)."""
        parts = section_id.split(".")
        try:
            if len(parts) >= 1:
                self.l1_counter = int(parts[0])
            if len(parts) >= 2:
                self.l2_counter = int(parts[1])
            else:
                self.l2_counter = 0
            if len(parts) >= 3:
                self.l3_counter = int(parts[2])
            else:
                self.l3_counter = 0
        except ValueError:
            # 비숫자가 섞인 경우 동기화 포기 — 다음 자동 번호는 기존 값 그대로 진행.
            pass

    # ---- shape handlers ---------------------------------------------

    def _handle_text_shape(self, section: Section, shape: Any) -> None:
        lines = extract_text_lines(shape)
        for text, marker in lines:
            if marker:
                section.blocks.append(
                    Block(type="list_item", text=text, marker=marker)
                )
            else:
                section.blocks.append(Block(type="paragraph", text=text))

    def _handle_table_shape(self, section: Section, shape: Any) -> None:
        headers, rows = extract_table_data(shape)
        if not headers and not rows:
            self.warnings.append(
                f"슬라이드 {section.id}: 빈 표 발견 — 무시함"
            )
            return

        self.tbl_counter += 1
        if not headers:
            width = max((len(r) for r in rows), default=1)
            headers = [f"col{i + 1}" for i in range(width)]
            self.warnings.append(
                f"표 {self.tbl_counter}: 헤더가 비어 있음 — 자동 헤더 사용"
            )
        if rows and any(len(r) != len(headers) for r in rows):
            self.warnings.append(
                f"표 {self.tbl_counter}: 일부 행이 헤더 길이와 다름"
            )

        tbl_id = _make_tbl_id(self.doc_id, self.tbl_counter)
        caption = f"Table {self.tbl_counter}: 슬라이드 {section.id} 표"
        tbl = Table(
            id=tbl_id,
            number=self.tbl_counter,
            caption=caption,
            section_ref=section.id,
            headers=headers,
            rows=rows,
        )
        self.tables.append(tbl)
        section.blocks.append(Block(type="table", ref=tbl_id))
        if tbl_id not in section.table_refs:
            section.table_refs.append(tbl_id)

    def _infer_caption_from_neighbor(self, shape: Any) -> str:
        """그림 도형에 어울리는 캡션을 alt-text / 인접 텍스트박스에서 추정.

        우선순위:
            1. python-pptx Picture 의 ``cNvPr@descr`` (alt-text) — 단순
               파일명 패턴이면 무시.
            2. 같은 슬라이드 안의 다른 텍스트박스 중, 가장 가까운 위치
               (그림 바로 아래 선호, 그 다음 위) 의 짧은 (≤120자) 첫 단락.
               ``Figure N:`` / ``Fig.`` / ``그림 N:`` 등 접두어가 있으면 우선.

        반환값은 prefix 없는 본문 또는 prefix 가 이미 붙은 전체 캡션이다.
        후보 없으면 빈 문자열.
        """
        # 1) alt-text 우선.
        alt = extract_picture_alt_text(shape)
        if alt:
            return alt

        # 2) 인접 텍스트박스 — 좌표 비교.
        # python-pptx Picture.\_parent 는 SlideShapes (Slide.shapes) 이며,
        # 직접 iterate 가능하다. iterable 이 아니면 ``.shapes`` 한 단계 더 시도.
        slide_obj = getattr(shape, "_parent", None)
        slide_shapes: Any = None
        if slide_obj is not None:
            shapes_attr = getattr(slide_obj, "shapes", None)
            if shapes_attr is not None and shapes_attr is not slide_obj:
                slide_shapes = shapes_attr
            else:
                slide_shapes = slide_obj
        if slide_shapes is None:
            return ""

        pic_left = getattr(shape, "left", None)
        pic_top = getattr(shape, "top", None)
        pic_width = getattr(shape, "width", None)
        pic_height = getattr(shape, "height", None)
        if None in (pic_left, pic_top, pic_width, pic_height):
            return ""

        pic_right = pic_left + pic_width
        pic_bottom = pic_top + pic_height

        # 후보 수집: (priority, distance, text)
        # priority: 0 = below + Figure-prefix, 1 = below 짧은 줄,
        #           2 = above + Figure-prefix, 3 = above 짧은 줄.
        candidates: list[tuple[int, float, str]] = []
        for other in slide_shapes:
            if other is shape:
                continue
            if not getattr(other, "has_text_frame", False):
                continue
            tf = getattr(other, "text_frame", None)
            if tf is None:
                continue
            text = (tf.text or "").strip()
            if not text:
                continue
            # 첫 줄만 캡션 후보로 본다.
            first_line = text.splitlines()[0].strip()
            if not first_line or len(first_line) > _CAPTION_MAX_LEN:
                continue

            o_left = getattr(other, "left", None)
            o_top = getattr(other, "top", None)
            o_width = getattr(other, "width", None)
            o_height = getattr(other, "height", None)
            if None in (o_left, o_top, o_width, o_height):
                continue
            o_bottom = o_top + o_height
            o_cx = o_left + o_width / 2

            # 가로 영역 겹침 — 그림 폭 안에 텍스트박스 중심이 들어와야 후보.
            if not (pic_left - pic_width * 0.2 <= o_cx <= pic_right + pic_width * 0.2):
                continue

            has_prefix = bool(_CAPTION_PREFIX_PATTERN.match(first_line))

            # 아래 위치 — 텍스트박스 top 이 그림 bottom 보다 같거나 큼.
            if o_top >= pic_bottom - pic_height * 0.05:
                dist = float(o_top - pic_bottom)
                candidates.append((0 if has_prefix else 1, dist, first_line))
                continue
            # 위 위치 — 텍스트박스 bottom 이 그림 top 보다 작거나 같음.
            if o_bottom <= pic_top + pic_height * 0.05:
                dist = float(pic_top - o_bottom)
                # "Figure N:" 접두어가 있으면 위쪽이라도 우선 (드물지만 발생).
                candidates.append((2 if has_prefix else 3, dist, first_line))
                continue

        if not candidates:
            return ""
        # priority 오름차순, 그 다음 거리 오름차순.
        candidates.sort(key=lambda t: (t[0], t[1]))
        return candidates[0][2]

    def _handle_picture_shape(self, section: Section, shape: Any) -> None:
        self.fig_counter += 1
        fig_id = _make_fig_id(self.doc_id, self.fig_counter)
        caption = f"Figure {self.fig_counter}: 슬라이드 {section.id} 이미지"

        # 캡션 자동 추정 휴리스틱 — alt-text 와 인접 텍스트박스 활용.
        if self.opts.infer_caption_from_neighbor:
            inferred = self._infer_caption_from_neighbor(shape)
            if inferred:
                # "Figure N:" 접두어가 이미 있으면 그대로, 없으면 prefix 자동 부여.
                if _CAPTION_PREFIX_PATTERN.match(inferred):
                    caption = inferred
                else:
                    caption = f"Figure {self.fig_counter}: {inferred}"

        fig = Figure(
            id=fig_id,
            number=self.fig_counter,
            caption=caption,
            section_ref=section.id,
        )
        self.figures.append(fig)
        section.blocks.append(Block(type="figure", ref=fig_id))
        if fig_id not in section.figure_refs:
            section.figure_refs.append(fig_id)

        # 첨부 일반화 (kind=figure) — Word 변환기와 동일.
        self.att_counter += 1
        att_id = _make_att_id(self.doc_id, self.att_counter)
        att = Attachment(
            id=att_id,
            number=self.att_counter,
            kind="figure",
            caption=caption,
            section_ref=section.id,
            extra={"figure_ref": fig_id},
        )
        self.attachments.append(att)

        # 바이너리 추출.
        if self.opts.extract_images:
            blob, ext = extract_picture_blob(shape)
            if blob:
                self._save_image(fig, att, blob, ext)
            else:
                self.warnings.append(
                    f"그림 {self.fig_counter}: 이미지 blob 추출 실패"
                )

    def _save_image(
        self, fig: Figure, att: Attachment, blob: bytes, ext: str
    ) -> None:
        out_root = Path(self.opts.output_dir) / self.doc_id
        try:
            out_root.mkdir(parents=True, exist_ok=True)
        except OSError as e:
            self.warnings.append(
                f"그림 출력 폴더 생성 실패: {out_root} ({e})"
            )
            return

        # 그림 본 파일.
        fig_name = f"F{fig.number:03d}.{ext}"
        fig_path = out_root / fig_name
        try:
            fig_path.write_bytes(blob)
        except OSError as e:
            self.warnings.append(
                f"그림 {fig.number}: 파일 쓰기 실패 — {fig_path} ({e})"
            )
            return
        fig.image_path = f"{self.doc_id}/{fig_name}"

        # 첨부 파일 경로 (POSIX-style).
        att_name = f"A{att.number:03d}.{ext}"
        att_path = out_root / att_name
        try:
            att_path.write_bytes(blob)
        except OSError as e:
            self.warnings.append(
                f"첨부 {att.number}: 파일 쓰기 실패 — {att_path} ({e})"
            )
            return
        att.file_name = att_name
        att.file_path = (Path(self.doc_id) / att_name).as_posix()
        att.mime_type = self._ext_to_mime(ext)
        try:
            att.size_bytes = att_path.stat().st_size
        except OSError:
            pass
        att.hash_sha256 = _sha256_of_file(att_path)

    @staticmethod
    def _ext_to_mime(ext: str) -> str:
        ext = ext.lower()
        return {
            "png": "image/png",
            "jpg": "image/jpeg",
            "jpeg": "image/jpeg",
            "gif": "image/gif",
            "bmp": "image/bmp",
            "tiff": "image/tiff",
            "wmf": "image/x-wmf",
            "emf": "image/x-emf",
        }.get(ext, "application/octet-stream")

    def _handle_chart_shape(self, section: Section, shape: Any) -> None:
        """차트 도형 처리.

        S7. 차트 데이터 추출:
            ``charts.extract_chart_table()`` 로 categories + series 값을
            ``Table`` 모델로 변환해 ``self.tables`` 에 등록한다. 동시에
            기존처럼 figure / attachment placeholder 도 만들어 본문 흐름을
            유지한다 (raw chart 자체는 이미지 추출이 어렵기 때문).
        """
        self.fig_counter += 1
        fig_id = _make_fig_id(self.doc_id, self.fig_counter)
        chart_title = extract_chart_title(shape)
        caption = f"Figure {self.fig_counter}: 차트"
        if chart_title:
            caption = f"Figure {self.fig_counter}: 차트 — {chart_title}"

        fig = Figure(
            id=fig_id,
            number=self.fig_counter,
            caption=caption,
            section_ref=section.id,
        )
        self.figures.append(fig)
        section.blocks.append(Block(type="figure", ref=fig_id))
        if fig_id not in section.figure_refs:
            section.figure_refs.append(fig_id)

        # 첨부 (kind=other — 차트는 임베드 부품 형태이므로 figure 가 아닌 일반 첨부로 둔다).
        self.att_counter += 1
        att_id = _make_att_id(self.doc_id, self.att_counter)

        # ---- S7. 차트 데이터 → Table -----------------------------------
        chart_table = None
        try:
            from .charts import extract_chart_table

            chart_table = extract_chart_table(shape)
        except Exception as exc:  # noqa: BLE001
            self.warnings.append(
                f"슬라이드 {section.id}: 차트 데이터 추출 실패 — {exc}"
            )

        chart_table_ref: str | None = None
        if chart_table is not None and not chart_table.is_empty():
            self.tbl_counter += 1
            chart_table_ref = _make_tbl_id(self.doc_id, self.tbl_counter)
            tbl_caption = (
                f"Table {self.tbl_counter}: 차트 데이터 — "
                f"{chart_title or chart_table.title or '제목없음'}"
            )
            tbl = Table(
                id=chart_table_ref,
                number=self.tbl_counter,
                caption=tbl_caption,
                section_ref=section.id,
                headers=list(chart_table.headers),
                rows=[list(r) for r in chart_table.rows],
            )
            self.tables.append(tbl)
            section.blocks.append(Block(type="table", ref=chart_table_ref))
            if chart_table_ref not in section.table_refs:
                section.table_refs.append(chart_table_ref)
        else:
            self.warnings.append(
                f"슬라이드 {section.id}: 차트 데이터 추출 결과 비어있음 (placeholder 만 등록)"
            )

        att = Attachment(
            id=att_id,
            number=self.att_counter,
            kind="chart",
            caption=caption,
            section_ref=section.id,
            extra={
                "figure_ref": fig_id,
                "chart_title": chart_title,
                "chart_type": (chart_table.chart_type if chart_table else ""),
                "chart_table_ref": chart_table_ref,
            },
        )
        self.attachments.append(att)

    # ---- 후처리: 본문 H2/H3 패턴 → sub-section 승격 ------------------

    def _extract_body_subsections(self) -> None:
        """슬라이드(level 1) section 의 blocks 중 본문 H2/H3 번호 패턴을
        만나는 paragraph/list_item 위치에서 새 sub-section 을 분기한다.

        실 데이터 패턴: 작성자가 슬라이드 제목은 그대로 두고 본문 첫 줄에
        "1.1 AI 도입 및 활용 현황" 같이 H2/H3 번호를 박아 둠. 이를 그대로
        평탄한 paragraph 로 두면 RAG 가 sub-주제를 인식하지 못한다.

        동작:
            - 매칭된 block 위치에서 새 sub-section 시작.
            - sub-section.id = 매칭 번호 ("1.1"), level = 점 개수 + 1.
            - sub-section.title = 번호 제거된 본문.
            - 다음 매칭 또는 슬라이드 끝까지의 blocks 를 sub-section.blocks 로
              이동 (단, 시작점 block 자체는 title 로 흡수되므로 제외).
            - 부모 슬라이드 section 의 blocks 에서는 이동된 blocks 를 제거.
            - id 충돌 시 (같은 슬라이드 안에서 동일 번호 두 번) 두 번째는
              자동 카운터로 폴백 + 경고.

        idempotent — 한 번 sub-section 이 만들어지면 그 children 안의
        blocks 는 다시 스캔하지 않는다.
        """
        if not self.section_root:
            return

        promoted = 0
        for parent in self.section_root:
            # children 안에 이미 부착된 sub-section 은 건드리지 않는다 — parent.blocks 만.
            blocks = parent.blocks
            # 매칭 위치 수집.
            matches: list[tuple[int, str, str]] = []  # (index, sub_id, sub_title)
            for i, blk in enumerate(blocks):
                btype = getattr(blk, "type", "")
                if btype not in ("paragraph", "list_item"):
                    continue
                text = (getattr(blk, "text", "") or "").strip()
                if not text:
                    continue
                sub_id, sub_title = infer_body_heading(text)
                if sub_id is None or not sub_title:
                    continue
                matches.append((i, sub_id, sub_title))

            if not matches:
                continue

            # 새 children 만들기 — 뒤에서부터 자르면 인덱스 보존.
            seen_ids: set[str] = set()
            # 부모의 leading blocks (첫 매칭 이전) 는 그대로 남는다.
            new_blocks = list(blocks[: matches[0][0]])
            new_children: list[Section] = []

            for k, (start, sub_id, sub_title) in enumerate(matches):
                end = matches[k + 1][0] if k + 1 < len(matches) else len(blocks)
                # 시작 block (제목으로 흡수) 다음부터 end 직전까지.
                child_blocks = list(blocks[start + 1 : end])

                # id 충돌 시 자동 카운터 폴백.
                if sub_id in seen_ids:
                    fallback_id = self._next_auto_id(level=2)
                    self.warnings.append(
                        f"슬라이드 {parent.id}: 본문 sub-section id 중복 "
                        f"({sub_id}) — 자동 id {fallback_id} 로 대체"
                    )
                    final_id = fallback_id
                    level = 2
                else:
                    final_id = sub_id
                    level = sub_id.count(".") + 1
                    seen_ids.add(sub_id)

                child = Section(
                    id=final_id,
                    level=level,
                    title=sub_title,
                    blocks=child_blocks,
                )
                # 자식에 속한 figure/table 참조도 부모에서 자식으로 옮긴다.
                self._migrate_refs_to_child(parent, child)
                new_children.append(child)
                promoted += 1

            parent.blocks = new_blocks
            # 기존 children (예: ``_open_section_for_slide`` 단계에서 붙은 것) 뒤에
            # 새 sub-section 을 추가.
            parent.children.extend(new_children)

        if promoted:
            self.warnings.append(
                f"본문 H2/H3 패턴 {promoted}건을 sub-section 으로 승격"
            )

    @staticmethod
    def _migrate_refs_to_child(parent: Section, child: Section) -> None:
        """child.blocks 안의 figure/table ref 를 parent 의 *_refs 에서 child 로 이동."""
        moved_figs: list[str] = []
        moved_tbls: list[str] = []
        for blk in child.blocks:
            ref = getattr(blk, "ref", None)
            if not ref:
                continue
            if getattr(blk, "type", "") == "figure":
                moved_figs.append(ref)
            elif getattr(blk, "type", "") == "table":
                moved_tbls.append(ref)

        for ref in moved_figs:
            if ref in parent.figure_refs:
                parent.figure_refs.remove(ref)
            if ref not in child.figure_refs:
                child.figure_refs.append(ref)
        for ref in moved_tbls:
            if ref in parent.table_refs:
                parent.table_refs.remove(ref)
            if ref not in child.table_refs:
                child.table_refs.append(ref)

    # ---- 후처리: 연속 동일 제목 그룹화 ------------------------------

    def _group_consecutive_duplicate_titles(self) -> None:
        """연속 동일 제목 슬라이드를 ``H1 시리즈`` 로 표기 — level 1 평탄 유지.

        실 PPT 의 흔한 패턴 — 작성자가 같은 제목으로 슬라이드 여러 장에 분할 작성.
        이는 의미상 H1 시리즈(분량 분리)이지 sub-주제(H2)가 아니다. 따라서 sub-grouping
        하지 않고 평탄한 level 1 을 유지하되, 제목에 ``(k/N)`` 위치 표기만 덧붙여
        같은 시리즈임을 식별 가능하게 한다.

            ✗ 이전: [parent (1/3)] children=[(2/3) level=2 id=1.1, (3/3) level=2 id=1.2]
                  → H1 H2 H2 트리, 의미상 sub-주제로 오해
            ✓ 지금: [(1/3) level=1 id=1, (2/3) level=1 id=2, (3/3) level=1 id=3]
                  → H1 H1 H1 시리즈, 형제 관계 유지

        ID 충돌 방지를 위해 자동 카운터가 부여한 unique id 는 그대로 보존하고,
        제목만 ``"X (k/N)"`` 형태로 표기한다.

        부작용 없음 — 변환 직후 1회만 호출된다 (idempotent).
        """
        if not self.section_root or len(self.section_root) < 2:
            return

        n = len(self.section_root)
        groups_made = 0
        i = 0
        while i < n:
            cur = self.section_root[i]
            cur_title = (cur.title or "").strip()
            if not cur_title:
                i += 1
                continue

            # 동일 제목 연속 길이.
            j = i + 1
            while j < n and (self.section_root[j].title or "").strip() == cur_title:
                j += 1
            run_len = j - i

            if run_len >= 2:
                # H1 시리즈 — 모든 슬라이드 level 1 유지, (k/N) 위치 표기만 추가.
                # id 는 자동 카운터가 이미 부여한 unique 값을 그대로 둔다.
                for k, idx in enumerate(range(i, j), start=1):
                    s = self.section_root[idx]
                    s.title = f"{cur_title} ({k}/{run_len})"
                groups_made += 1

            i = j

        if groups_made:
            self.warnings.append(
                f"연속 동일 제목 {groups_made}건을 H1 시리즈로 표기 (level 1 평탄 유지, (k/N) 위치만 표기)"
            )

    # ---- summary 폴백 ------------------------------------------------

    def _extract_summary_fallback(self, core: Any) -> str:
        """meta.summary 후보 텍스트를 다단계 폴백으로 추출.

        1) core_properties.subject → 가장 권위 있는 신호.
        2) 표지 다음 슬라이드(또는 1슬라이드뿐이면 슬라이드 1)의 paragraph/list_item
           텍스트를 모아 약 250자 이내로 자름.

        둘 다 없으면 빈 문자열. 어떤 슬라이드 본문이라도 있으면 빈 summary 보다는
        낫다 (RAG 검색 가중치 향상).
        """
        # 1) core_properties.subject 우선.
        if core is not None:
            subject = (getattr(core, "subject", "") or "").strip()
            if subject and len(subject) >= 10:
                return subject[:300]

        # 2) 본문 슬라이드 텍스트 추출.
        if not self.section_root:
            return ""
        # 슬라이드가 2개 이상이면 표지(1번)를 건너뛰고 2번부터.
        skip_first = len(self.section_root) > 1
        parts: list[str] = []
        total_len = 0
        for i, sec in enumerate(self.section_root):
            if skip_first and i == 0:
                continue
            for blk in sec.blocks:
                if getattr(blk, "type", "") not in ("paragraph", "list_item"):
                    continue
                t = (getattr(blk, "text", "") or "").strip()
                if not t or len(t) < 5:
                    continue
                # speaker notes 마커 제외.
                if t in ("[Speaker Notes]",):
                    continue
                parts.append(t)
                total_len += len(t)
                if total_len >= 250:
                    break
            if total_len >= 250:
                break

        if not parts:
            return ""
        joined = " · ".join(parts)
        return joined[:300] + ("…" if len(joined) > 300 else "")

    # ---- meta --------------------------------------------------------

    def _build_meta(
        self, prs: PresentationType, pptx_path: Path
    ) -> dict[str, Any]:
        try:
            core = prs.core_properties
        except Exception:  # noqa: BLE001
            core = None

        title = ""
        author = ""
        created_str = ""
        modified_str = ""
        if core is not None:
            title = (getattr(core, "title", "") or "").strip()
            author = (getattr(core, "author", "") or "").strip()
            created = getattr(core, "created", None)
            modified = getattr(core, "modified", None)
            now = datetime.now(tz=timezone.utc)
            created_str = (created or now).strftime("%Y-%m-%d")
            modified_str = (modified or now).strftime("%Y-%m-%d")

        if not title:
            title = pptx_path.stem
        if not created_str:
            created_str = datetime.now(tz=timezone.utc).strftime("%Y-%m-%d")
        if not modified_str:
            modified_str = created_str

        # summary 폴백 (default ON)
        summary = ""
        if self.opts.extract_summary:
            summary = self._extract_summary_fallback(core)

        meta: dict[str, Any] = {
            "doc_id": self.doc_id,
            "title": title,
            "source_format": "pptx",
            "source_file": pptx_path.name,
            "doc_type": "slide",
            "created": created_str,
            "modified": modified_str,
            "author": author,
            "department": f"{self.opts.division}-{self.opts.team}",
            "version": "1.0",
            "tags": list(self.opts.tags),
            "summary": summary,
        }
        if self.opts.agents:
            meta["agent_scope"] = list(self.opts.agents)

        # Migration 0007: agent-discovery 자동 기본값.
        _apply_agent_discovery_defaults(
            meta,
            overrides=None,  # PPT 변환기는 meta override 채널이 따로 없음.
            data_type_name="PowerPoint 슬라이드",
            title=str(title),
            tags=meta["tags"],
            section_count=len(self.section_root),
            table_count=len(self.tables),
            figure_count=len(self.figures),
        )

        if not meta["tags"]:
            self.warnings.append("tags 미지정 → 빈 배열")
        if not meta["summary"]:
            self.warnings.append("summary 미지정 → 빈 문자열")

        return meta


# ---------------------------------------------------------------------------
# Output writing
# ---------------------------------------------------------------------------
def write_output(
    result: ConversionResult,
    output_dir: Path,
) -> tuple[Path, Path]:
    """JSON 과 경고 로그 저장. (json_path, log_path) 반환."""
    import json

    output_dir.mkdir(parents=True, exist_ok=True)

    doc_id = result.meta["doc_id"]
    json_path = output_dir / f"{doc_id}.json"
    log_path = output_dir / f"{doc_id}.warnings.log"

    json_path.write_text(
        json.dumps(result.to_dict(), ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    if result.warnings:
        log_path.write_text(
            "\n".join(f"[WARN] {w}" for w in result.warnings),
            encoding="utf-8",
        )
    elif log_path.exists():
        log_path.unlink()

    return json_path, log_path


__all__ = [
    "PptxConverter",
    "PptxConverterOptions",
    "write_output",
]
