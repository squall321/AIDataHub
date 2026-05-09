"""표준 PowerPoint(.pptx) 예제 생성 스크립트.

생성되는 sample_presentation.pptx 가 시연하는 원칙
(ppt_to_json_conversion_rules.md):
- 제목 placeholder 사용 (4.1 절 — 제목 추출 우선순위 1)
- 본문 텍스트박스 (자연스러운 불릿 — 변환기는 paragraph/list_item 으로 추출)
- 표 도형 (7.1 절 — 표 추출)
- 그림 + 캡션 텍스트박스 (7.4 절 — 캡션 작성 표준)
- 발표자 노트 (6장 — 노트 추출)

CAE/IGA 도메인 테마 — iga_guide.docx 와 일관.

실행 방법
--------
python _generate_ppt.py [출력경로]
기본 출력: 같은 폴더의 sample_presentation.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def build_sample_presentation(out_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 워크북/프레젠테이션 빌트인 속성 (core_properties)
    # → PPT 변환기가 meta.title / meta.author 의 폴백으로 사용한다.
    cp = prs.core_properties
    cp.title = "KooRemapper IGA 변환 결과 (표준 예제)"
    cp.author = "CAE팀"
    cp.subject = "IGA 검증 발표자료"
    cp.keywords = "IGA, NURBS, KooRemapper"
    cp.comments = (
        "표준 PPT 작성 예제 — 슬라이드 제목 placeholder, 텍스트박스 라벨, "
        "표 도형, 그림 placeholder + 캡션, 발표자 노트를 모두 시연한다."
    )

    # ── Slide 1: 제목 슬라이드 ────────────────────────────────────
    title_layout = prs.slide_layouts[0]  # Title Slide
    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = "KooRemapper IGA 변환 결과 (표준 예제)"
    s1.placeholders[1].text = (
        "CAE팀 · 2026-05-08\nstandard PPT 작성 예제"
    )

    # ── Slide 2: "1. 개요" + 불릿 본문 ────────────────────────────
    blank_layout = prs.slide_layouts[5]  # Title Only
    s2 = prs.slides.add_slide(blank_layout)
    s2.shapes.title.text = "1. 개요"
    tb = s2.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(8.5), Inches(5.0)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    # 자연스러운 불릿 본문 — 다중 단락 + level 들여쓰기로 list_item 으로 추출됨.
    bullets = [
        ("KooRemapper 의 IGA 변환 결과를 검증한다.", 0),
        ("Trimmed NURBS Volume 방식으로 FE mesh 를 IGA 입력으로 변환.", 0),
        ("LS-DYNA 재해석으로 FEM 대비 응력 분포 차이를 1% 이내로 확인.", 0),
        ("브라켓 부품 50,000 노드 기준 — 변환 8.7 초, 검증 0.3 초.", 1),
    ]
    first = True
    for text, lvl in bullets:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        para.text = text
        para.level = lvl
        para.font.size = Pt(20)
        first = False

    # ── Slide 3: "2. 작동 원리" + 표 ──────────────────────────────
    s3 = prs.slides.add_slide(blank_layout)
    s3.shapes.title.text = "2. 작동 원리"
    rows, cols = 3, 3
    table_shape = s3.shapes.add_table(
        rows, cols,
        Inches(0.7), Inches(1.5),
        Inches(8.5), Inches(2.5),
    )
    tbl = table_shape.table
    headers = ["단계", "입력", "출력"]
    data = [
        ("Bbox 계산", "FE mesh", "min/max XYZ"),
        ("NURBS 박스 생성", "bbox + offset", "8개 제어점"),
    ]
    for j, h in enumerate(headers):
        tbl.cell(0, j).text = h
    for i, row in enumerate(data, start=1):
        for j, v in enumerate(row):
            tbl.cell(i, j).text = v

    # ── Slide 4: "3. 검증" + 그림 영역 + 캡션 + 발표자 노트 ───────
    s4 = prs.slides.add_slide(blank_layout)
    s4.shapes.title.text = "3. 검증"

    # 그림 영역 (단순 도형 — 실제 이미지가 없을 때 대체용)
    placeholder = s4.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(2.5), Inches(1.5),
        Inches(5.0), Inches(3.0),
    )
    placeholder.text = "NURBS 박스 다이어그램"

    # 캡션 텍스트박스 — 그림 바로 아래 (Figure N: ... 패턴)
    cap_box = s4.shapes.add_textbox(
        Inches(2.5), Inches(4.7),
        Inches(5.0), Inches(0.5),
    )
    cap_box.text_frame.text = "Figure 1: NURBS 박스가 FE mesh 를 4 mm 오프셋으로 둘러싼다."
    cap_box.text_frame.paragraphs[0].font.size = Pt(14)
    cap_box.text_frame.paragraphs[0].font.italic = True

    # 발표자 노트
    notes = s4.notes_slide.notes_text_frame
    notes.text = (
        "이 슬라이드는 검증 결과를 시연한다. 다음 슬라이드에서 정량 비교를 진행."
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def main() -> int:
    out = (
        Path(sys.argv[1])
        if len(sys.argv) > 1
        else Path(__file__).resolve().parent / "sample_presentation.pptx"
    )
    p = build_sample_presentation(out)
    print(f"[OK] PPT 예제 생성: {p}  ({p.stat().st_size:,} bytes)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
