"""PPT(.pptx) → JSON 변환기 단위 테스트.

python-pptx 의 ``Presentation()`` 으로 작은 가상 프레젠테이션을 만들고,
``PptxConverter.convert()`` 가 올바른 ConversionResult 를 만들어내는지 검증한다.

테스트 대상 슬라이드:
    1) 제목 + 글머리 본문 + 슬라이드 노트
    2) 번호 패턴 제목("1.2 작동원리") + 본문
    3) 제목 + 표 (3x3) — 표 처리 검증
    4) 제목 + 그림 1장 — 첨부/이미지 추출 검증
"""
from __future__ import annotations

import io
import struct
import zlib
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from ppt_converter.core import (
    PptxConverter,
    PptxConverterOptions,
    write_output,
)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------
def _make_minimal_png() -> bytes:
    """1x1 단색 PNG 바이트 생성 (외부 파일 의존 없음)."""

    def _chunk(tag: bytes, data: bytes) -> bytes:
        return (
            struct.pack(">I", len(data))
            + tag
            + data
            + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
        )

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = _chunk(
        b"IHDR",
        struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0),
    )
    raw = b"\x00\xff\x00\x00"  # filter byte + RGB pixel
    idat = _chunk(b"IDAT", zlib.compress(raw))
    iend = _chunk(b"IEND", b"")
    return sig + ihdr + idat + iend


@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    """테스트용 .pptx 를 생성하여 파일 경로 반환."""
    prs = Presentation()
    blank = prs.slide_layouts[5]  # "Title Only" — placeholder 0 = title
    title_layout = prs.slide_layouts[1]  # "Title and Content"

    # ---- Slide 1: title + bullets + notes ---------------------------
    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = "도입 및 배경"
    body = s1.placeholders[1]
    tf = body.text_frame
    tf.text = "첫 번째 항목"
    p2 = tf.add_paragraph()
    p2.text = "두 번째 항목"
    p3 = tf.add_paragraph()
    p3.text = "세 번째 항목"
    s1.notes_slide.notes_text_frame.text = "발표 시 보충 설명: 배경을 강조"

    # ---- Slide 2: numbered title -----------------------------------
    s2 = prs.slides.add_slide(title_layout)
    s2.shapes.title.text = "1.2 작동원리"
    s2.placeholders[1].text_frame.text = "원리 설명 본문"

    # ---- Slide 3: title + table ------------------------------------
    s3 = prs.slides.add_slide(blank)
    s3.shapes.title.text = "성능 비교"
    rows, cols = 3, 3
    left = top = Inches(1.0)
    width = Inches(6.0)
    height = Inches(2.0)
    tbl_shape = s3.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    headers = ["항목", "방법A", "방법B"]
    data = [
        ["속도", "100", "150"],
        ["정확도", "0.9", "0.95"],
    ]
    for c, h in enumerate(headers):
        tbl.cell(0, c).text = h
    for r, row in enumerate(data, start=1):
        for c, v in enumerate(row):
            tbl.cell(r, c).text = v

    # ---- Slide 4: title + picture ----------------------------------
    s4 = prs.slides.add_slide(blank)
    s4.shapes.title.text = "다이어그램"
    img_stream = io.BytesIO(_make_minimal_png())
    s4.shapes.add_picture(
        img_stream, Inches(1.0), Inches(1.0), Inches(2.0), Inches(2.0)
    )

    out = tmp_path / "sample.pptx"
    prs.save(str(out))
    return out


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------
def _build_options(tmp_path: Path) -> PptxConverterOptions:
    return PptxConverterOptions(
        team="HE",
        group="CAE",
        year=2026,
        seq=1,
        output_dir=tmp_path / "out",
        tags=["테스트", "PPT"],
        agents=["iga-analyst"],
    )


def test_basic_conversion_produces_sections(sample_pptx: Path, tmp_path: Path) -> None:
    converter = PptxConverter(_build_options(tmp_path))
    result = converter.convert(sample_pptx)

    # 4개의 슬라이드 → 최소 4개의 섹션 (level 2 의 1.2 는 부모가 없어 root 폴백 가능).
    assert len(result.sections) >= 3, (
        f"섹션이 너무 적음: {[(s.id, s.title) for s in result.sections]}"
    )

    # toc 합산 (children 포함) 으로 4개 보장.
    toc = result.build_toc()
    assert len(toc) == 4, f"toc 항목 수 불일치: {toc}"

    titles = [item["title"] for item in toc]
    assert "도입 및 배경" in titles
    assert "작동원리" in titles
    assert "성능 비교" in titles
    assert "다이어그램" in titles


def test_numbered_title_extracts_section_id(
    sample_pptx: Path, tmp_path: Path
) -> None:
    converter = PptxConverter(_build_options(tmp_path))
    result = converter.convert(sample_pptx)

    toc = result.build_toc()
    matches = [item for item in toc if item["title"] == "작동원리"]
    assert matches, f"'작동원리' 섹션이 toc 에 없음: {toc}"
    assert matches[0]["id"] == "1.2", (
        f"제목에서 섹션 번호 1.2 가 추출되지 않음: {matches[0]}"
    )


def test_table_is_extracted_with_ref(
    sample_pptx: Path, tmp_path: Path
) -> None:
    converter = PptxConverter(_build_options(tmp_path))
    result = converter.convert(sample_pptx)

    assert len(result.tables) == 1, (
        f"표가 1개 추출되어야 함: {len(result.tables)}"
    )
    tbl = result.tables[0]
    assert tbl.headers == ["항목", "방법A", "방법B"]
    assert len(tbl.rows) == 2

    # 표가 속한 섹션이 ref 를 가져야 함.
    table_section = None
    for section in result.sections:
        if tbl.id in section.table_refs:
            table_section = section
            break
    assert table_section is not None, "표 ref 가 섹션에 등록되지 않음"
    assert table_section.title == "성능 비교"


def test_picture_creates_figure_and_attachment(
    sample_pptx: Path, tmp_path: Path
) -> None:
    converter = PptxConverter(_build_options(tmp_path))
    result = converter.convert(sample_pptx)

    assert len(result.figures) == 1, (
        f"그림이 1개 추출되어야 함: {len(result.figures)}"
    )
    fig = result.figures[0]
    assert fig.caption.startswith("Figure 1:"), (
        f"caption 누락 — {fig.caption!r}"
    )

    # 첨부 (kind=figure) 1건이 함께 등록되어야 함.
    figure_attachments = [a for a in result.attachments if a.kind == "figure"]
    assert len(figure_attachments) == 1
    att = figure_attachments[0]
    assert att.caption == fig.caption
    # POSIX-style 상대경로 (cross-platform) 확인.
    assert att.file_path is not None
    assert "/" in att.file_path
    assert "\\" not in att.file_path

    # 실제 파일이 디스크에 저장되어야 함.
    out_root = (tmp_path / "out") / result.meta["doc_id"]
    saved_files = list(out_root.glob("*"))
    assert saved_files, f"이미지 파일이 저장되지 않음: {out_root}"


def test_speaker_notes_appear_with_marker(
    sample_pptx: Path, tmp_path: Path
) -> None:
    converter = PptxConverter(_build_options(tmp_path))
    result = converter.convert(sample_pptx)

    intro_section = next(
        (s for s in result.sections if s.title == "도입 및 배경"), None
    )
    assert intro_section is not None

    block_texts = [b.text for b in intro_section.blocks if b.text]
    assert "[Speaker Notes]" in block_texts, (
        f"speaker notes 마커가 없음: {block_texts}"
    )
    assert any("배경을 강조" in t for t in block_texts), (
        f"speaker notes 본문이 누락: {block_texts}"
    )


def test_meta_fields_set(sample_pptx: Path, tmp_path: Path) -> None:
    converter = PptxConverter(_build_options(tmp_path))
    result = converter.convert(sample_pptx)

    meta = result.meta
    assert meta["doc_id"] == "DOC-HE-CAE-2026-000001"
    assert meta["source_format"] == "pptx"
    assert meta["doc_type"] == "slide"
    assert meta["department"] == "HE-CAE"
    assert meta["tags"] == ["테스트", "PPT"]
    assert meta["agent_scope"] == ["iga-analyst"]


def test_write_output_creates_files(sample_pptx: Path, tmp_path: Path) -> None:
    opts = _build_options(tmp_path)
    converter = PptxConverter(opts)
    result = converter.convert(sample_pptx)

    json_path, log_path = write_output(result, opts.output_dir)
    assert json_path.exists()
    assert json_path.suffix == ".json"

    import json as _json

    payload = _json.loads(json_path.read_text(encoding="utf-8"))
    assert payload["meta"]["doc_id"] == result.meta["doc_id"]
    assert payload["schema_version"] == "1.0"
    assert isinstance(payload["sections"], list)
    assert isinstance(payload["tables"], list)
    assert isinstance(payload["figures"], list)
    assert isinstance(payload["attachments"], list)
