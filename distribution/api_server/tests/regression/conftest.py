"""실제 변환기 산출물을 in-memory 로 만들어 회귀 테스트에 공급한다.

루트 ``tests/conftest.py`` 가 이미 PG → SQLite 어댑터를 등록하고 ``test_engine``
/ ``test_session_maker`` / ``db_client`` 를 제공하므로, 본 conftest 는 그 위에
실제 .docx / .xlsx / .pptx / .md / .pdf 파일을 생성하는 픽스처와, ingest →
API 호출 헬퍼만 추가한다.

각 fixture 는 ``tmp_path`` 위에 파일을 만들어 반환하므로 호출 측에서 직접
``converter`` 모듈을 호출해 변환할 수 있다. 외부 바이너리 의존은 모두
optional(dev) 패키지로, 미설치 시 해당 fixture 는 skip 된다.
"""
from __future__ import annotations

import struct
import textwrap
import zlib
from pathlib import Path
from typing import Any

import pytest


# ---------------------------------------------------------------------------
# 1x1 PNG (binary attachment 용)
# ---------------------------------------------------------------------------
def _make_minimal_png() -> bytes:
    """완전히 self-contained 한 1x1 PNG 바이트.

    외부 픽스처 디렉터리에 의존하지 않고 매 테스트 실행마다 동일 결과를
    만들어내야 하므로 파이썬으로 직접 인코딩한다 (test_ppt_converter 와
    동일 알고리즘).
    """

    def _chunk(tag: bytes, data: bytes) -> bytes:
        return (
            struct.pack(">I", len(data))
            + tag
            + data
            + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
        )

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = _chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
    raw = b"\x00\xff\x00\x00"
    idat = _chunk(b"IDAT", zlib.compress(raw))
    iend = _chunk(b"IEND", b"")
    return sig + ihdr + idat + iend


@pytest.fixture(scope="session")
def png_bytes() -> bytes:
    return _make_minimal_png()


# ---------------------------------------------------------------------------
# DOCX fixture — 제목 1/2/3 + 단락 + 표 + 이미지 + 코드블록
# ---------------------------------------------------------------------------
@pytest.fixture()
def sample_docx(tmp_path: Path, png_bytes: bytes) -> Path:
    """h1+h2+h3 + paragraph + table + image + code block 포함 .docx 파일."""
    docx = pytest.importorskip("docx")
    doc = docx.Document()

    doc.add_heading("실시간 회귀 테스트 문서", level=1)
    doc.add_paragraph("이 문서는 회귀 스위트가 동적으로 생성합니다.")

    doc.add_heading("배경 및 동기", level=2)
    doc.add_paragraph(
        "Migration 0006 의 ``classification`` 등 메타가 인제스트 후에도 "
        "올바르게 채워지는지 검증합니다."
    )

    doc.add_heading("세부 절차", level=3)
    doc.add_paragraph("1) 변환 → 2) 정규화 → 3) DB 적재 → 4) API 조회")

    # 표
    table = doc.add_table(rows=3, cols=2)
    table.style = "Table Grid"
    table.cell(0, 0).text = "항목"
    table.cell(0, 1).text = "값"
    table.cell(1, 0).text = "ID"
    table.cell(1, 1).text = "DOC-HE-CAE-2026-100001"
    table.cell(2, 0).text = "분류"
    table.cell(2, 1).text = "internal"

    # 이미지
    img_path = tmp_path / "_sample.png"
    img_path.write_bytes(png_bytes)
    doc.add_picture(str(img_path))
    doc.add_paragraph("그림 1. 회귀 테스트 마커 이미지")

    # 코드블록 (paragraph w/ Courier 스타일이 없어도 텍스트로만 OK)
    code_p = doc.add_paragraph()
    run = code_p.add_run("for i in range(3):\n    print(i)")
    run.font.name = "Consolas"

    out = tmp_path / "regression_word.docx"
    doc.save(str(out))
    return out


# ---------------------------------------------------------------------------
# XLSX fixture — 2 시트 + _META + _GLOSSARY + 단위 + 병합셀
# ---------------------------------------------------------------------------
@pytest.fixture()
def sample_xlsx(tmp_path: Path) -> Path:
    """2 개 시트 + _META + _GLOSSARY 포함 .xlsx 파일.

    구조:
        Sheet "측정데이터": 헤더 + 단위 행 + 데이터 행, 1행 병합 ("배터리 시험 결과").
        Sheet "재시도": 두 번째 표.
        Sheet "_META": key/value.
        Sheet "_GLOSSARY": term/definition.
    """
    openpyxl = pytest.importorskip("openpyxl")
    wb = openpyxl.Workbook()
    default = wb.active
    wb.remove(default)

    # ---- 측정데이터 -----------------------------------------------------
    ws = wb.create_sheet(title="측정데이터")
    # 1행: 병합셀 헤더
    ws.cell(row=1, column=1, value="배터리 시험 결과")
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=3)
    # 2행: 컬럼명
    ws.append(["time", "force", "strain"])
    # 3행: 단위
    ws.append(["[ms]", "[N]", "[%]"])
    # 4-6행: 값
    ws.append([0.0, 0.0, 0.0])
    ws.append([0.1, 12.5, 0.02])
    ws.append([0.2, 25.0, 0.05])

    # ---- 재시도 ---------------------------------------------------------
    ws2 = wb.create_sheet(title="재시도")
    ws2.append(["sample", "value"])
    ws2.append(["A", 1])
    ws2.append(["B", 2])

    # ---- _META ---------------------------------------------------------
    meta = wb.create_sheet(title="_META")
    meta.append(["key", "value"])
    meta.append(["title", "배터리 측정 데이터"])
    meta.append(["author", "qa-bot"])
    meta.append(["classification", "internal"])
    meta.append(["domain", "battery"])

    # ---- _GLOSSARY ----------------------------------------------------
    g = wb.create_sheet(title="_GLOSSARY")
    g.append(["term", "definition"])
    g.append(["force", "노일 관통 충격 하중"])
    g.append(["strain", "변형률 (%)"])

    out = tmp_path / "regression_excel.xlsx"
    wb.save(out)
    wb.close()
    return out


# ---------------------------------------------------------------------------
# PPTX fixture — 제목 + 본문 + 표 + 발표 노트 + (의사) 차트 슬라이드
# ---------------------------------------------------------------------------
@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx.Presentation()
    title_layout = prs.slide_layouts[1]  # Title and Content
    blank = prs.slide_layouts[5]

    # Slide 1 — title + body + notes
    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = "회귀 테스트 발표"
    body = s1.placeholders[1]
    body.text_frame.text = "주요 변경사항 요약"
    p2 = body.text_frame.add_paragraph()
    p2.text = "Migration 0006 백필 결과"
    s1.notes_slide.notes_text_frame.text = (
        "발표 노트: 회귀 스위트가 자동 생성한 슬라이드입니다."
    )

    # Slide 2 — title + 3x3 table
    s2 = prs.slides.add_slide(blank)
    s2.shapes.title.text = "성능 비교"
    rows, cols = 3, 3
    tbl = s2.shapes.add_table(
        rows, cols, Inches(1.0), Inches(1.5), Inches(6.0), Inches(2.0)
    ).table
    headers = ["지표", "이전", "이후"]
    data = [["속도", "100", "150"], ["정확도", "0.90", "0.95"]]
    for c, h in enumerate(headers):
        tbl.cell(0, c).text = h
    for r, row in enumerate(data, start=1):
        for c, v in enumerate(row):
            tbl.cell(r, c).text = v

    # Slide 3 — title + dummy "chart" represented as a text body
    # (실제 차트 임베드는 python-pptx 가 별도 API 를 요구하므로, ppt_converter
    # 가 chart 부재 환경에서도 동작하는지 확인하기 위한 텍스트 placeholder).
    s3 = prs.slides.add_slide(title_layout)
    s3.shapes.title.text = "차트 (텍스트 표현)"
    s3.placeholders[1].text_frame.text = "차트 데이터: 100/150"

    out = tmp_path / "regression_ppt.pptx"
    prs.save(str(out))
    return out


# ---------------------------------------------------------------------------
# MD fixture — front matter + h1-h3 + table + image
# ---------------------------------------------------------------------------
@pytest.fixture()
def sample_md(tmp_path: Path) -> Path:
    text = textwrap.dedent(
        """\
        ---
        title: 회귀 테스트 마크다운
        author: qa-bot
        classification: internal
        domain: testing
        ---

        # 도입

        본 문서는 회귀 테스트가 자동으로 생성합니다.

        ## 절차

        1. 변환
        2. 적재
        3. 검증

        ### 세부 (h3)

        세부 단계 설명.

        | 단계 | 결과 |
        |------|------|
        | 변환 | OK   |
        | 적재 | OK   |

        ![캡션 예시](./image.png)
        """
    )
    out = tmp_path / "regression_md.md"
    out.write_text(text, encoding="utf-8")
    return out


# ---------------------------------------------------------------------------
# PDF fixture — outline 헤딩 + paragraph + table
# ---------------------------------------------------------------------------
@pytest.fixture()
def sample_pdf(tmp_path: Path) -> Path:
    rl = pytest.importorskip("reportlab")  # noqa: F841
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import (
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table as RLTable,
        TableStyle,
    )

    out = tmp_path / "regression_pdf.pdf"
    styles = getSampleStyleSheet()

    story: list[Any] = []
    story.append(Paragraph("<b>회귀 테스트 PDF</b>", styles["Title"]))
    story.append(Spacer(1, 12))
    story.append(Paragraph("<b>1. 도입</b>", styles["Heading1"]))
    story.append(
        Paragraph("이 PDF 는 회귀 테스트가 자동으로 만든 문서입니다.", styles["BodyText"])
    )
    story.append(Spacer(1, 8))
    story.append(Paragraph("<b>2. 데이터</b>", styles["Heading1"]))
    table_data = [
        ["항목", "값"],
        ["속도", "150"],
        ["정확도", "0.95"],
    ]
    rl_tbl = RLTable(table_data, hAlign="LEFT")
    rl_tbl.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.black),
            ]
        )
    )
    story.append(rl_tbl)

    doc = SimpleDocTemplate(str(out), pagesize=A4)
    doc.build(story)
    return out


# ---------------------------------------------------------------------------
# Helper: ingest a payload via writer + commit, return WriteResult.
# ---------------------------------------------------------------------------
@pytest.fixture()
def regression_ingest(test_session_maker):
    """``payload(dict) → WriteResult`` 함수 fixture.

    테스트는 변환기로 payload 를 만들고 이 헬퍼를 호출해 DB 에 적재한다.
    """
    from api.ingest.db_writer import write_record
    from api.ingest.normalizer import normalize

    async def _ingest(payload: dict[str, Any]):
        record_in = normalize(payload)
        async with test_session_maker() as session:
            result = await write_record(session, record_in)
            await session.commit()
        return result, record_in

    return _ingest
