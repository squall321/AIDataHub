"""S7. PPT chart data extraction.

생성된 .pptx 에 작은 막대 차트를 삽입하고, 변환기 출력의 ``tables``
배열에 차트 데이터가 등록되는지 검증한다.

python-pptx 가 없으면 모듈 단위 skip.
"""
from __future__ import annotations

import io
from pathlib import Path

import pytest


pytest.importorskip("pptx")


def _make_pptx_with_chart(tmp_path: Path) -> Path:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank)

    # 제목
    title = slide.shapes.title
    if title is not None:
        title.text = "차트 슬라이드"

    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3"]
    chart_data.add_series("매출", (10.0, 20.0, 30.0))
    chart_data.add_series("비용", (5.0, 8.0, 12.0))

    x, y, cx, cy = Inches(1), Inches(2), Inches(6), Inches(4)
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )

    p = tmp_path / "with_chart.pptx"
    prs.save(p)
    return p


def test_extract_chart_table_unit(tmp_path) -> None:
    """charts.extract_chart_table 단위 동작."""
    from pptx import Presentation

    from ppt_converter.charts import extract_chart_table
    from ppt_converter.parser import is_chart_shape

    pptx_path = _make_pptx_with_chart(tmp_path)
    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]

    chart_shape = None
    for shape in slide.shapes:
        if is_chart_shape(shape):
            chart_shape = shape
            break
    assert chart_shape is not None

    ct = extract_chart_table(chart_shape)
    assert ct is not None
    assert ct.headers[:1] == ["category"]
    assert "매출" in ct.headers
    assert "비용" in ct.headers
    # rows: Q1/Q2/Q3.
    assert len(ct.rows) == 3
    cats = [r[0] for r in ct.rows]
    assert cats == ["Q1", "Q2", "Q3"]
    # 매출 column.
    sales_idx = ct.headers.index("매출")
    sales_values = [r[sales_idx] for r in ct.rows]
    assert sales_values == [10.0, 20.0, 30.0]


def test_pptx_converter_emits_chart_table(tmp_path) -> None:
    """PptxConverter.convert() 가 차트를 ``tables[]`` 에 등록한다."""
    from ppt_converter.core import PptxConverter, PptxConverterOptions

    pptx_path = _make_pptx_with_chart(tmp_path)
    opts = PptxConverterOptions(
        team="HE",
        group="CAE",
        year=2026,
        seq=99,
        output_dir=tmp_path / "out",
    )
    conv = PptxConverter(opts)
    result = conv.convert(pptx_path)

    # 표 1 개 이상 — 차트 데이터에서 유래.
    assert len(result.tables) >= 1
    chart_table = result.tables[0]
    assert chart_table.headers[0] == "category"
    assert any(h in chart_table.headers for h in ("매출", "비용"))
    # 본문 흐름에 type=table block 이 들어갔는지.
    found_block = False
    for s in result.sections:
        for b in s.blocks:
            if getattr(b, "type", None) == "table":
                found_block = True
        for c in s.children:
            for b in c.blocks:
                if getattr(b, "type", None) == "table":
                    found_block = True
    assert found_block, "table block missing in slide section"
