"""python-pptx 헬퍼 — slide / shape 단위 추출 유틸.

이 모듈은 python-pptx 객체에 대한 얕은 wrapper 역할만 한다 (값 추출, ID 추정 등).
실제 섹션 트리 / blocks 조립은 ``ppt_converter.core`` 가 담당한다.
"""
from __future__ import annotations

import re
from dataclasses import dataclass
from typing import Any, Iterator

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation as PresentationType
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture
from pptx.slide import Slide

# "1. 제목", "1.2 작동원리", "1.2.3 세부 절차" 등 — Word 변환기와 동일 패턴.
SECTION_NUM_PATTERN = re.compile(r"^(\d+(?:\.\d+){0,2})[\.\)]?\s+(.*)$")

# 본문 H2/H3 휴리스틱 — 슬라이드 제목과 달리 적어도 한 점이 있어야 한다
# ("1." 단독 단락은 단순 목록 번호일 가능성이 높아 제외).
BODY_HEADING_PATTERN = re.compile(r"^(\d+\.\d+(?:\.\d+)?)[\.\)]?\s+(.+)$")


# ---------------------------------------------------------------------------
# Slide-level helpers
# ---------------------------------------------------------------------------
@dataclass
class SlideInfo:
    """변환 단계에서 사용할 슬라이드 메타."""

    index: int  # 1-based
    slide: Slide


def iter_slides_with_index(prs: PresentationType) -> Iterator[SlideInfo]:
    """프레젠테이션의 모든 슬라이드를 1-based 인덱스와 함께 yield."""
    for i, slide in enumerate(prs.slides, start=1):
        yield SlideInfo(index=i, slide=slide)


def extract_slide_title(slide: Slide) -> str:
    """슬라이드 제목 추출.

    우선순위:
        1) ``slide.shapes.title`` placeholder 의 텍스트
        2) placeholder 중 ``ph.type`` 이 TITLE / CTR_TITLE 인 것
        3) 첫 번째 텍스트 프레임의 첫 줄 (제목 placeholder 가 없는 슬라이드)
    """
    title_shape = slide.shapes.title
    if title_shape is not None:
        text = (title_shape.text_frame.text or "").strip()
        if text:
            return text

    # 폴백 1: placeholder 중 title 타입
    for ph in slide.placeholders:
        ph_type = getattr(ph.placeholder_format, "type", None)
        type_name = getattr(ph_type, "name", "") if ph_type is not None else ""
        if type_name in {"TITLE", "CTR_TITLE"}:
            text = (ph.text_frame.text or "").strip()
            if text:
                return text

    # 폴백 2: 첫 텍스트 프레임의 첫 줄
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = (shape.text_frame.text or "").strip()
            if text:
                return text.splitlines()[0].strip()

    return ""


def extract_speaker_notes(slide: Slide) -> str:
    """슬라이드 노트 텍스트 추출. 노트가 없으면 빈 문자열."""
    if not slide.has_notes_slide:
        return ""
    notes_slide = slide.notes_slide
    if notes_slide is None:
        return ""
    tf = notes_slide.notes_text_frame
    if tf is None:
        return ""
    return (tf.text or "").strip()


# ---------------------------------------------------------------------------
# Shape-level extractors
# ---------------------------------------------------------------------------
def _is_title_placeholder(shape: BaseShape) -> bool:
    """이 도형이 슬라이드 제목 placeholder 인지 판단."""
    if not getattr(shape, "is_placeholder", False):
        return False
    ph_format = getattr(shape, "placeholder_format", None)
    if ph_format is None:
        return False
    ph_type = getattr(ph_format, "type", None)
    type_name = getattr(ph_type, "name", "") if ph_type is not None else ""
    return type_name in {"TITLE", "CTR_TITLE", "SUBTITLE"}


def iter_body_shapes(slide: Slide) -> Iterator[BaseShape]:
    """제목/SUBTITLE placeholder 를 제외한 본문 도형들을 등장 순서대로 yield.

    ``shape.shape_type`` 이 GROUP 이면 자식 도형까지 평탄화한다.
    """
    for shape in slide.shapes:
        if _is_title_placeholder(shape):
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_group(shape)
        else:
            yield shape


def _iter_group(group_shape: BaseShape) -> Iterator[BaseShape]:
    """GroupShape 내 자식들을 평탄화하여 yield (group 중첩도 재귀)."""
    children = getattr(group_shape, "shapes", None)
    if children is None:
        return
    for child in children:
        if child.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_group(child)
        else:
            yield child


def extract_text_lines(shape: BaseShape) -> list[tuple[str, str | None]]:
    """텍스트 프레임에서 (text, marker) 쌍 리스트 추출.

    marker 는 list_item 처리용 — 단락의 ``level`` (들여쓰기 깊이) 에 따라
    "•", "  ◦", "    ▪" 등을 부여한다. 단락 텍스트가 비어 있으면 결과에 포함하지 않는다.
    """
    if not getattr(shape, "has_text_frame", False):
        return []
    tf = shape.text_frame
    out: list[tuple[str, str | None]] = []
    paragraphs = list(getattr(tf, "paragraphs", []))

    # 단일 단락이고 모든 단락이 평문(들여쓰기 0) 일 때는 list_item 화 하지 않는다.
    multi_or_indented = len(paragraphs) > 1 or any(
        getattr(p, "level", 0) > 0 for p in paragraphs
    )

    for p in paragraphs:
        text = "".join(run.text for run in p.runs).strip()
        if not text:
            continue
        level = getattr(p, "level", 0) or 0
        if multi_or_indented:
            indent = "  " * level
            marker = f"{indent}•"
            out.append((text, marker))
        else:
            out.append((text, None))
    return out


def extract_table_data(graphic_frame: GraphicFrame) -> tuple[list[str], list[list[Any]]]:
    """GraphicFrame.has_table → headers / rows.

    첫 행을 헤더로 가정. 빈 셀은 None.
    """
    if not getattr(graphic_frame, "has_table", False):
        return [], []
    tbl = graphic_frame.table
    rows_iter = list(tbl.rows)
    if not rows_iter:
        return [], []

    headers: list[str] = []
    for cell in rows_iter[0].cells:
        text = (cell.text or "").strip()
        headers.append(text or f"col{len(headers) + 1}")

    data_rows: list[list[Any]] = []
    for row in rows_iter[1:]:
        row_vals: list[Any] = []
        for cell in row.cells:
            t = (cell.text or "").strip()
            row_vals.append(_coerce_cell_value(t))
        data_rows.append(row_vals)
    return headers, data_rows


def _coerce_cell_value(text: str) -> Any:
    """Word 변환기와 동일한 셀 값 정규화 — 빈 문자열은 None,
    정수/실수 패턴은 숫자로 변환, 그 외는 그대로 문자열 보존."""
    if text == "":
        return None
    if re.fullmatch(r"-?\d+", text):
        try:
            return int(text)
        except ValueError:
            return text
    if re.fullmatch(r"-?\d+\.\d+", text):
        try:
            return float(text)
        except ValueError:
            return text
    return text


def is_picture_shape(shape: BaseShape) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.PICTURE


def is_chart_shape(shape: BaseShape) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.CHART


def is_table_shape(shape: BaseShape) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.TABLE or getattr(shape, "has_table", False)


def is_text_shape(shape: BaseShape) -> bool:
    """텍스트가 있는 도형 (placeholder 본문, 텍스트박스 등)."""
    return getattr(shape, "has_text_frame", False)


# Picture cNvPr@descr (alt-text) 추출용 네임스페이스 토큰.
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

# alt-text 가 단순 파일명일 때는 caption 으로 쓰지 않는다 — 이미지 라이브러리가
# 자동으로 채우는 placeholder 인 경우가 대부분.
_TRIVIAL_ALT_PATTERN = re.compile(
    r"^(image|picture|그림|사진|photo|img|pic)\d*\.[a-z]{2,5}$",
    re.IGNORECASE,
)


def extract_picture_alt_text(shape: BaseShape) -> str:
    """Picture 도형의 ``cNvPr@descr`` (alt-text) 를 추출.

    파일명 패턴 ("image.png", "사진1.jpg" 등) 은 의미 있는 캡션이 아니므로
    빈 문자열을 반환한다.
    """
    el = getattr(shape, "_element", None)
    if el is None:
        el = getattr(shape, "element", None)
    if el is None:
        return ""
    # p:nvPicPr/p:cNvPr — descr 또는 title 속성.
    for cnv in el.iter(f"{{{_NS_P}}}cNvPr"):
        descr = (cnv.get("descr") or "").strip()
        if descr and not _TRIVIAL_ALT_PATTERN.fullmatch(descr):
            return descr
        title = (cnv.get("title") or "").strip()
        if title and not _TRIVIAL_ALT_PATTERN.fullmatch(title):
            return title
        # 첫 cNvPr 만 검사 (Picture 도형은 1개).
        break
    return ""


def extract_picture_blob(shape: Picture) -> tuple[bytes, str]:
    """Picture 도형의 바이너리와 확장자 추출. 추출 실패 시 (b"", "")."""
    image = getattr(shape, "image", None)
    if image is None:
        return b"", ""
    blob = getattr(image, "blob", b"") or b""
    ext = (getattr(image, "ext", "") or "").lstrip(".").lower()
    if not ext:
        ct = (getattr(image, "content_type", "") or "").lower()
        if "/" in ct:
            ext = ct.split("/", 1)[1].split(";", 1)[0].strip()
            if ext == "jpeg":
                ext = "jpg"
    return blob, ext or "png"


def extract_chart_title(shape: Shape) -> str:
    """차트 제목 추출 — 없으면 빈 문자열. python-pptx 의 chart 객체 사용."""
    chart = getattr(shape, "chart", None)
    if chart is None:
        return ""
    if not getattr(chart, "has_title", False):
        return ""
    title = getattr(chart, "chart_title", None)
    if title is None:
        return ""
    tf = getattr(title, "text_frame", None)
    if tf is None:
        return ""
    return (tf.text or "").strip()


# ---------------------------------------------------------------------------
# Section ID inference
# ---------------------------------------------------------------------------
def infer_section_id_from_title(title: str) -> tuple[str | None, str]:
    """슬라이드 제목에서 ``"1.2 작동원리"`` 같은 번호를 분리.

    Returns:
        (section_id_or_None, cleaned_title)
    """
    if not title:
        return None, ""
    m = SECTION_NUM_PATTERN.match(title.strip())
    if not m:
        return None, title.strip()
    return m.group(1), m.group(2).strip()


def infer_body_heading(text: str) -> tuple[str | None, str]:
    """본문 paragraph/list_item 텍스트에서 H2/H3 번호 패턴을 추출.

    슬라이드 제목용 ``infer_section_id_from_title`` 와 달리, 적어도 한 점이
    있는 번호만 매칭한다 (``1.1``, ``1.2.3``). ``1`` 또는 ``1.`` 단독은
    단순 목록 번호일 가능성이 높으므로 제외한다.

    Returns:
        (section_id_or_None, cleaned_title)
    """
    if not text:
        return None, ""
    m = BODY_HEADING_PATTERN.match(text.strip())
    if not m:
        return None, text.strip()
    return m.group(1), m.group(2).strip()


__all__ = [
    "SECTION_NUM_PATTERN",
    "BODY_HEADING_PATTERN",
    "SlideInfo",
    "iter_slides_with_index",
    "extract_slide_title",
    "extract_speaker_notes",
    "iter_body_shapes",
    "extract_text_lines",
    "extract_table_data",
    "is_picture_shape",
    "is_chart_shape",
    "is_table_shape",
    "is_text_shape",
    "extract_picture_blob",
    "extract_picture_alt_text",
    "extract_chart_title",
    "infer_section_id_from_title",
    "infer_body_heading",
]
